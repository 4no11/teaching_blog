import os
import json
import requests
from flask import Blueprint, request, jsonify, Response
from config import Config
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
from werkzeug.utils import secure_filename

ALLOWED_EXTENSIONS = {'ppt', 'pptx', 'doc', 'docx', 'pdf'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_ppt(filepath):
    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        return f"PPT解析错误: {str(e)}"

def extract_text_from_docx(filepath):
    try:
        doc = Document(filepath)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return '\n'.join(text)
    except Exception as e:
        return f"Word解析错误: {str(e)}"

def extract_text_from_pdf(filepath):
    try:
        reader = PdfReader(filepath)
        text = []
        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)
        return '\n'.join(text)
    except Exception as e:
        return f"PDF解析错误: {str(e)}"

def extract_text_from_file(filepath, ext):
    if ext in ['ppt', 'pptx']:
        return extract_text_from_ppt(filepath)
    elif ext in ['doc', 'docx']:
        return extract_text_from_docx(filepath)
    elif ext == 'pdf':
        return extract_text_from_pdf(filepath)
    return "不支持的文件格式"

ai_bp = Blueprint('ai', __name__)

@ai_bp.route('/upload-courseware', methods=['POST'])
def upload_courseware():
    if 'file' not in request.files:
        return jsonify({'error': '没有上传文件'})
    
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'error': '没有选择文件'})
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        
        if '.' not in filename:
            return jsonify({'error': '文件没有扩展名，请上传PPT、Word或PDF文件'})
        
        ext = filename.rsplit('.', 1)[1].lower()
        
        upload_folder = os.path.join(os.getcwd(), 'static', 'uploads')
        os.makedirs(upload_folder, exist_ok=True)
        
        filepath = os.path.join(upload_folder, filename)
        file.save(filepath)
        
        try:
            text = extract_text_from_file(filepath, ext)
            
            if text.startswith(('PPT解析错误', 'Word解析错误', 'PDF解析错误', '不支持')):
                return jsonify({'error': text})
            
            os.remove(filepath)
            
            return jsonify({
                'success': True,
                'content': text,
                'filename': filename
            })
        except Exception as e:
            if os.path.exists(filepath):
                os.remove(filepath)
            return jsonify({'error': f'文件处理失败: {str(e)}'})
    
    return jsonify({'error': '不支持的文件格式，请上传PPT、Word或PDF文件'})

def call_ai_api(messages, max_tokens=1000, retries=2):
    if not Config.AI_API_KEY:
        return "AI配置错误：请设置 AI_API_KEY 环境变量"

    headers = {
        'Authorization': f'Bearer {Config.AI_API_KEY}',
        'Content-Type': 'application/json'
    }
    
    data = {
        'model': Config.AI_MODEL,
        'messages': messages,
        'max_tokens': max_tokens,
        'temperature': 0.7
    }
    
    for attempt in range(retries):
        try:
            response = requests.post(
                f'{Config.AI_BASE_URL}/chat/completions',
                headers=headers,
                json=data,
                timeout=60
            )
            response.raise_for_status()
            result = response.json()
            
            if 'choices' in result:
                return result['choices'][0]['message']['content']
            else:
                return f"AI调用失败: {result}"
        except requests.exceptions.Timeout:
            if attempt < retries - 1:
                continue
            return "AI调用失败: 请求超时，请稍后重试"
        except Exception as e:
            return f"AI调用失败: {str(e)}"

def stream_ai_api(messages, max_tokens=1000):
    if not Config.AI_API_KEY:
        yield "data: {\"error\": \"AI配置错误：请设置 AI_API_KEY 环境变量\"}\n\n"
        return

    headers = {
        'Authorization': f'Bearer {Config.AI_API_KEY}',
        'Content-Type': 'application/json'
    }
    
    data = {
        'model': Config.AI_MODEL,
        'messages': messages,
        'max_tokens': max_tokens,
        'temperature': 0.7,
        'stream': True
    }
    
    try:
        response = requests.post(
            f'{Config.AI_BASE_URL}/chat/completions',
            headers=headers,
            json=data,
            timeout=120,
            stream=True
        )
        response.raise_for_status()
        
        for line in response.iter_lines():
            if line:
                line = line.decode('utf-8')
                if line.startswith('data: ') and line != 'data: [DONE]':
                    yield f"{line}\n\n"
                    
    except Exception as e:
        yield "data: {\"error\": \"" + str(e) + "\"}\n\n"

# 会话存储，用于管理上下文
conversation_store = {}

@ai_bp.route('/chat', methods=['POST'])
def chat():
    data = request.get_json()
    user_message = data.get('message', '')
    session_id = data.get('session_id', 'default')
    
    # 初始化会话
    if session_id not in conversation_store:
        conversation_store[session_id] = [
            {'role': 'system', 'content': '你是一个教育助手，可以帮助教师回答教学相关问题，提供教学建议和资源推荐。'}
        ]
    
    # 添加用户消息
    conversation_store[session_id].append({'role': 'user', 'content': user_message})
    
    # 限制会话长度，防止过长
    if len(conversation_store[session_id]) > 10:
        # 保留系统消息和最近的8条消息
        conversation_store[session_id] = [conversation_store[session_id][0]] + conversation_store[session_id][-8:]
    
    response = call_ai_api(conversation_store[session_id])
    
    # 添加AI回复到会话
    conversation_store[session_id].append({'role': 'assistant', 'content': response})
    
    return jsonify({'response': response, 'session_id': session_id})

@ai_bp.route('/chat/stream', methods=['POST'])
def chat_stream():
    data = request.get_json()
    user_message = data.get('message', '')
    session_id = data.get('session_id', 'default')
    
    # 初始化会话
    if session_id not in conversation_store:
        conversation_store[session_id] = [
            {'role': 'system', 'content': '你是一个教育助手，可以帮助教师回答教学相关问题，提供教学建议和资源推荐。'}
        ]
    
    # 添加用户消息
    conversation_store[session_id].append({'role': 'user', 'content': user_message})
    
    # 限制会话长度
    if len(conversation_store[session_id]) > 10:
        conversation_store[session_id] = [conversation_store[session_id][0]] + conversation_store[session_id][-8:]
    
    def generate():
        full_response = ""
        for chunk in stream_ai_api(conversation_store[session_id]):
            yield chunk
            try:
                chunk_data = json.loads(chunk.replace('data: ', '').strip())
                if 'choices' in chunk_data:
                    content = chunk_data['choices'][0].get('delta', {}).get('content', '')
                    full_response += content
                elif 'error' in chunk_data:
                    full_response = chunk_data['error']
            except:
                pass
        
        # 完成后保存完整回复到会话
        if full_response and not full_response.startswith('AI调用失败') and not full_response.startswith('AI配置错误'):
            conversation_store[session_id].append({'role': 'assistant', 'content': full_response})
    
    return Response(
        generate(),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@ai_bp.route('/chat/clear', methods=['POST'])
def clear_chat():
    data = request.get_json()
    session_id = data.get('session_id', 'default')
    
    if session_id in conversation_store:
        del conversation_store[session_id]
    
    return jsonify({'success': True})

@ai_bp.route('/summarize', methods=['POST'])
def summarize():
    data = request.get_json()
    content = data.get('content', '')
    
    messages = [
        {'role': 'system', 'content': '你是一个文本摘要助手，请将输入的文本简洁地总结成一段话。'},
        {'role': 'user', 'content': f'请总结以下内容：\n{content}'}
    ]
    
    summary = call_ai_api(messages, max_tokens=300)
    return jsonify({'summary': summary})

@ai_bp.route('/improve', methods=['POST'])
def improve():
    data = request.get_json()
    content = data.get('content', '')
    
    messages = [
        {'role': 'system', 'content': '你是一个写作改进助手，请改进输入的文本，使其更加清晰、流畅，并提高可读性。保持原意不变。'},
        {'role': 'user', 'content': f'请改进以下文本：\n{content}'}
    ]
    
    improved = call_ai_api(messages, max_tokens=2000)
    return jsonify({'improved': improved})

@ai_bp.route('/generate-title', methods=['POST'])
def generate_title():
    data = request.get_json()
    content = data.get('content', '')
    
    messages = [
        {'role': 'system', 'content': '你是一个标题生成助手，根据内容生成一个吸引人的标题，最多50个字符。'},
        {'role': 'user', 'content': f'根据以下内容生成标题：\n{content}'}
    ]
    
    title = call_ai_api(messages, max_tokens=100)
    return jsonify({'title': title.strip('"')})

@ai_bp.route('/suggest-tags', methods=['POST'])
def suggest_tags():
    data = request.get_json()
    content = data.get('content', '')
    
    messages = [
        {'role': 'system', 'content': '你是一个标签推荐助手，根据内容推荐5个相关的标签，用逗号分隔。'},
        {'role': 'user', 'content': f'根据以下内容推荐标签：\n{content}'}
    ]
    
    tags = call_ai_api(messages, max_tokens=100)
    return jsonify({'tags': tags})

@ai_bp.route('/lesson-plan', methods=['POST'])
def generate_lesson_plan():
    data = request.get_json()
    topic = data.get('topic', '')
    grade = data.get('grade', '')
    duration = data.get('duration', '45')
    
    if not topic:
        return jsonify({'error': '请输入教学知识点'})
    
    messages = [
        {'role': 'system', 'content': '''你是一个专业的教学设计专家。请简洁地生成教案，包含以下要点：

## 教案信息
- 课题：[名称]
- 适用年级：[年级]
- 课时：[时长]

## 教学目标
1. 知识目标：
2. 能力目标：
3. 情感目标：

## 教学重难点
- 重点：
- 难点：

## 教学过程（约500字）
1. 导入（5分钟）：
2. 新课讲授（25分钟）：
3. 课堂活动（10分钟）：
4. 总结（5分钟）：

## 作业设计
（1-2项作业）'''},
        {'role': 'user', 'content': f'''请为以下内容生成简洁教案：
知识点：{topic}
年级：{grade}
时长：{duration}分钟'''}
    ]
    
    plan = call_ai_api(messages, max_tokens=800)
    return jsonify({'plan': plan})

@ai_bp.route('/quiz', methods=['POST'])
def generate_quiz():
    data = request.get_json()
    topic = data.get('topic', '')
    question_type = data.get('type', 'all')
    difficulty = data.get('difficulty', 'medium')
    count = data.get('count', 5)
    
    if not topic:
        return jsonify({'error': '请输入教学知识点'})
    
    type_desc = {
        'choice': '选择题',
        'fill': '填空题',
        'short': '简答题',
        'all': '选择题、填空题、简答题'
    }
    
    messages = [
        {'role': 'system', 'content': '''你是一个出题专家。请根据知识点生成练习题，格式：

## 练习题

### 选择题
1. [题目]
   A. 
   B. 
   C. 
   D. 
   **答案：A** **解析：** 

### 填空题
1. [题目]
   **答案：** 

### 简答题
1. [题目]
   **答案：** 

题目简洁准确，答案明确。'''},
        {'role': 'user', 'content': f'''生成{count}道{difficulty}难度的{type_desc.get(question_type, '各类题型')}：
知识点：{topic}'''}
    ]
    
    quiz = call_ai_api(messages, max_tokens=800)
    return jsonify({'quiz': quiz})

@ai_bp.route('/extract-knowledge', methods=['POST'])
def extract_knowledge():
    data = request.get_json()
    content = data.get('content', '')
    title = data.get('title', '')
    
    if not content:
        return jsonify({'error': '请输入课件内容'})
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的教育知识提取专家。请从课件内容中提取核心知识点和重点内容：

## 知识点提取结果

### 核心知识点
（列出5-10个核心知识点，按重要性排序，每个知识点用简短的条目形式）

### 重点内容
（列出3-5个最重要的重点内容，需要特别强调的部分）

### 知识结构
（分析知识点之间的逻辑关系，用简短的结构描述）

### 教学重点
（指出教学过程中需要重点讲解的内容）

### 难点内容
（指出学生可能难以理解的内容）

请确保提取的知识点准确、全面，结构清晰。'''
        },
        {
            'role': 'user',
            'content': f'''课件标题：{title}
课件内容：
{content}'''
        }
    ]
    
    result = call_ai_api(messages, max_tokens=1000)
    return jsonify({'knowledge': result})

@ai_bp.route('/extract-knowledge-stream', methods=['POST'])
def extract_knowledge_stream():
    """流式提取知识点"""
    data = request.get_json()
    content = data.get('content', '')
    title = data.get('title', '')

    if not content:
        return Response("data: {\"error\": \"请输入课件内容\"}\n\n", mimetype='text/event-stream')

    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的教育知识提取专家。请从课件内容中提取核心知识点和重点内容：

## 知识点提取结果

### 核心知识点
（列出5-10个核心知识点，按重要性排序，每个知识点用简短的条目形式）

### 重点内容
（列出3-5个最重要的重点内容，需要特别强调的部分）

### 知识结构
（分析知识点之间的逻辑关系，用简短的结构描述）

### 教学重点
（指出教学过程中需要重点讲解的内容）

### 难点内容
（指出学生可能难以理解的内容）

请确保提取的知识点准确、全面，结构清晰。'''
        },
        {
            'role': 'user',
            'content': f'''课件标题：{title}
课件内容：
{content}'''
        }
    ]

    def generate():
        for chunk in stream_ai_api(messages, max_tokens=1000):
            try:
                if chunk.startswith('data: '):
                    json_str = chunk[6:].strip()
                    if json_str and json_str != '[DONE]':
                        parsed = json.loads(json_str)
                        choices = parsed.get('choices', [])
                        if choices:
                            delta = choices[0].get('delta', {})
                            content = delta.get('content', '')
                            if content:
                                yield f"data: {{\"content\": {json.dumps(content, ensure_ascii=False)}}}\n\n"
            except:
                pass

    return Response(generate(), mimetype='text/event-stream')

@ai_bp.route('/courseware-analyze', methods=['POST'])
def analyze_courseware():
    data = request.get_json()
    content = data.get('content', '')
    title = data.get('title', '')
    
    if not content:
        return jsonify({'error': '请输入课件内容'})
    
    messages = [
        {'role': 'system', 'content': '''你是一个教学设计专家。请分析课件内容，从以下方面提出优化建议：

## 课件分析报告

### 一、内容完整性
（评估知识点的覆盖程度和深度）

### 二、逻辑结构
（评估内容组织是否清晰、层次是否分明）

### 三、可读性评估
（评估表述是否易懂、专业术语是否过多）

### 四、优化建议
1. 内容改进建议
2. 结构优化建议
3. 表达优化建议

### 五、特色亮点
（指出课件的优点）

请给出具体、可操作的改进建议。'''},
        {'role': 'user', 'content': f'''课件标题：{title}
课件内容：
{content}'''}
    ]
    
    result = call_ai_api(messages, max_tokens=1000)
    return jsonify({'analysis': result})

@ai_bp.route('/courseware-analyze-stream', methods=['POST'])
def analyze_courseware_stream():
    """流式课件分析"""
    data = request.get_json()
    content = data.get('content', '')
    title = data.get('title', '')

    if not content:
        return Response("data: {\"error\": \"请输入课件内容\"}\n\n", mimetype='text/event-stream')

    messages = [
        {'role': 'system', 'content': '''你是一个教学设计专家。请分析课件内容，从以下方面提出优化建议：

## 课件分析报告

### 一、内容完整性
（评估知识点的覆盖程度和深度）

### 二、逻辑结构
（评估内容组织是否清晰、层次是否分明）

### 三、可读性评估
（评估表述是否易懂、专业术语是否过多）

### 四、优化建议
1. 内容改进建议
2. 结构优化建议
3. 表达优化建议

### 五、特色亮点
（指出课件的优点）

请给出具体、可操作的改进建议。'''},
        {'role': 'user', 'content': f'''课件标题：{title}
课件内容：
{content}'''}
    ]

    def generate():
        for chunk in stream_ai_api(messages, max_tokens=1000):
            try:
                if chunk.startswith('data: '):
                    json_str = chunk[6:].strip()
                    if json_str and json_str != '[DONE]':
                        parsed = json.loads(json_str)
                        choices = parsed.get('choices', [])
                        if choices:
                            delta = choices[0].get('delta', {})
                            content = delta.get('content', '')
                            if content:
                                yield f"data: {{\"content\": {json.dumps(content, ensure_ascii=False)}}}\n\n"
            except:
                pass

    return Response(generate(), mimetype='text/event-stream')

@ai_bp.route('/courseware-summary', methods=['POST'])
def generate_summary_objectives():
    data = request.get_json()
    content = data.get('content', '')
    
    if not content:
        return jsonify({'error': '请输入课件内容'})
    
    messages = [
        {'role': 'system', 'content': '''你是一个教学设计专家。请根据课件内容生成摘要和学习目标：

## 内容摘要
（用2-3段话概括课件核心内容，150字左右）

## 学习目标

### 知识与技能目标
1. 

### 过程与方法目标
1. 

### 情感态度与价值观目标
1. 

## 核心知识点清单
（列出5-8个核心知识点，用简短条目形式）

要求表述简洁、明确，便于教师参考和学生理解。'''},
        {'role': 'user', 'content': f'''请分析以下课件内容：
{content}'''}
    ]
    
    result = call_ai_api(messages, max_tokens=800)
    return jsonify({'summary': result})

@ai_bp.route('/courseware-summary-stream', methods=['POST'])
def generate_summary_objectives_stream():
    """流式生成摘要目标"""
    data = request.get_json()
    content = data.get('content', '')

    if not content:
        return Response("data: {\"error\": \"请输入课件内容\"}\n\n", mimetype='text/event-stream')

    messages = [
        {'role': 'system', 'content': '''你是一个教学设计专家。请根据课件内容生成摘要和学习目标：

## 内容摘要
（用2-3段话概括课件核心内容，150字左右）

## 学习目标

### 知识与技能目标
1. 

### 过程与方法目标
1. 

### 情感态度与价值观目标
1. 

## 核心知识点清单
（列出5-8个核心知识点，用简短条目形式）

要求表述简洁、明确，便于教师参考和学生理解。'''},
        {'role': 'user', 'content': f'''请分析以下课件内容：
{content}'''}
    ]

    def generate():
        for chunk in stream_ai_api(messages, max_tokens=800):
            try:
                if chunk.startswith('data: '):
                    json_str = chunk[6:].strip()
                    if json_str and json_str != '[DONE]':
                        parsed = json.loads(json_str)
                        choices = parsed.get('choices', [])
                        if choices:
                            delta = choices[0].get('delta', {})
                            content = delta.get('content', '')
                            if content:
                                yield f"data: {{\"content\": {json.dumps(content, ensure_ascii=False)}}}\n\n"
            except:
                pass

    return Response(generate(), mimetype='text/event-stream')

@ai_bp.route('/courseware-layout', methods=['POST'])
def suggest_layout():
    data = request.get_json()
    content = data.get('content', '')
    page_count = data.get('page_count', 10)
    
    if not content:
        return jsonify({'error': '请输入课件内容'})
    
    messages = [
        {'role': 'system', 'content': '''你是一个课件排版专家。请根据内容提供智能排版建议：

## 排版建议方案

### 一、整体结构规划
建议将内容分为X个部分，每部分X页

### 二、页面分配
- 导入/封面：1页
- 知识讲解：X页
- 示例/案例：X页
- 练习/互动：X页
- 总结/作业：1页

### 三、布局建议
1. 每页信息量控制（建议不超过X个要点）
2. 字体大小建议
3. 配色方案建议
4. 图表/配图建议

### 四、重点内容突出
（建议哪些内容需要特别强调、如何突出）

### 五、视觉层次
（建议标题、正文、注释的层级关系）

请给出实用、具体的排版指南。'''},
        {'role': 'user', 'content': f'''课件共约{page_count}页，内容如下：
{content}'''}
    ]
    
    result = call_ai_api(messages, max_tokens=800)
    return jsonify({'layout': result})

@ai_bp.route('/courseware-layout-stream', methods=['POST'])
def suggest_layout_stream():
    """流式生成排版建议"""
    data = request.get_json()
    content = data.get('content', '')
    page_count = data.get('page_count', 10)

    if not content:
        return Response("data: {\"error\": \"请输入课件内容\"}\n\n", mimetype='text/event-stream')

    messages = [
        {'role': 'system', 'content': '''你是一个课件排版专家。请根据内容提供智能排版建议：

## 排版建议方案

### 一、整体结构规划
建议将内容分为X个部分，每部分X页

### 二、页面分配
- 导入/封面：1页
- 知识讲解：X页
- 示例/案例：X页
- 练习/互动：X页
- 总结/作业：1页

### 三、布局建议
1. 每页信息量控制（建议不超过X个要点）
2. 字体大小建议
3. 配色方案建议
4. 图表/配图建议

### 四、重点内容突出
（建议哪些内容需要特别强调、如何突出）

### 五、视觉层次
（建议标题、正文、注释的层级关系）

请给出实用、具体的排版指南。'''},
        {'role': 'user', 'content': f'''课件共约{page_count}页，内容如下：
{content}'''}
    ]

    def generate():
        for chunk in stream_ai_api(messages, max_tokens=800):
            try:
                if chunk.startswith('data: '):
                    json_str = chunk[6:].strip()
                    if json_str and json_str != '[DONE]':
                        parsed = json.loads(json_str)
                        choices = parsed.get('choices', [])
                        if choices:
                            delta = choices[0].get('delta', {})
                            content = delta.get('content', '')
                            if content:
                                yield f"data: {{\"content\": {json.dumps(content, ensure_ascii=False)}}}\n\n"
            except:
                pass

    return Response(generate(), mimetype='text/event-stream')

@ai_bp.route('/teaching-reflection', methods=['POST'])
def generate_reflection():
    data = request.get_json()
    process = data.get('process', '')
    subject = data.get('subject', '')
    grade = data.get('grade', '')
    
    if not process:
        return jsonify({'error': '请输入教学过程记录'})
    
    messages = [
        {'role': 'system', 'content': '''你是一个资深教学专家。请根据教学过程记录，生成教学反思要点和改进建议：

## 教学反思报告

### 一、教学亮点
（回顾本节课的成功之处）

### 二、不足与反思
1. 教学设计层面：
2. 教学过程层面：
3. 学生反馈层面：

### 三、改进建议
1. 短期改进（下次课可调整）：
2. 中期改进（学期内可优化）：
3. 长期发展方向：

### 四、教学策略推荐
（针对当前教学内容，推荐3-5种有效的教学策略）

### 五、备选方案
（如时间充裕可以增加的环节）

请给出真诚、务实的反思建议。'''},
        {'role': 'user', 'content': f'''教学基本信息:
- 学科:{subject}
- 年级:{grade}

教学过程记录:
{process}'''}
    ]
    
    result = call_ai_api(messages, max_tokens=1000)
    return jsonify({'reflection': result})

@ai_bp.route('/teaching-reflection-stream', methods=['POST'])
def generate_reflection_stream():
    """流式生成教学反思"""
    data = request.get_json()
    process = data.get('process', '')
    subject = data.get('subject', '')
    grade = data.get('grade', '')

    if not process:
        return Response("data: {\"error\": \"请输入教学过程记录\"}\n\n", mimetype='text/event-stream')

    messages = [
        {'role': 'system', 'content': '''你是一个资深教学专家。请根据教学过程记录，生成教学反思要点和改进建议：

## 教学反思报告

### 一、教学亮点
（回顾本节课的成功之处）

### 二、不足与反思
1. 教学设计层面：
2. 教学过程层面：
3. 学生反馈层面：

### 三、改进建议
1. 短期改进（下次课可调整）：
2. 中期改进（学期内可优化）：
3. 长期发展方向：

### 四、教学策略推荐
（针对当前教学内容，推荐3-5种有效的教学策略）

### 五、备选方案
（如时间充裕可以增加的环节）

请给出真诚、务实的反思建议。'''},
        {'role': 'user', 'content': f'''教学基本信息:
- 学科:{subject}
- 年级:{grade}

教学过程记录:
{process}'''}
    ]

    def generate():
        for chunk in stream_ai_api(messages, max_tokens=1000):
            try:
                if chunk.startswith('data: '):
                    json_str = chunk[6:].strip()
                    if json_str and json_str != '[DONE]':
                        parsed = json.loads(json_str)
                        choices = parsed.get('choices', [])
                        if choices:
                            delta = choices[0].get('delta', {})
                            content = delta.get('content', '')
                            if content:
                                yield f"data: {{\"content\": {json.dumps(content, ensure_ascii=False)}}}\n\n"
            except:
                pass

    return Response(generate(), mimetype='text/event-stream')

# ========== 一键备课功能相关接口 ==========

@ai_bp.route('/generate-lesson-plan-stream', methods=['POST'])
def generate_lesson_plan_stream():
    """流式生成教案"""
    data = request.get_json()
    stage = data.get('stage', '高中')
    requirements = data.get('requirements', '')
    file_content = data.get('fileContent', '')  # 接收文件内容
    file_name = data.get('fileName', '')
    
    # 构建用户消息
    user_message = f'请为{stage}阶段设计一份完整教案。\n\n'
    
    # 如果有教材内容，添加到prompt中
    if file_content and len(file_content.strip()) > 0:
        user_message += f'''【上传的教材内容】（文件名：{file_name}）
请严格基于以下教材内容来设计教案，确保教案内容与教材紧密相关：

{file_content}

---
'''
    
    user_message += f'''用户自定义要求：
{requirements if requirements else "无特殊要求，请根据通用标准设计"}

请开始生成教案：'''
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的教学设计专家，擅长为不同学科和年级设计高质量教案。
请严格按照以下格式生成完整的教案内容：

# 教案标题

## 一、教学目标

### 1. 知识与技能目标
（列出具体知识点目标）

### 2. 过程与方法目标
（列出能力培养目标）

### 3. 情感态度与价值观目标
（列出情感培养目标）

## 二、教学分析

### 1. 教材分析
（分析教材地位和作用）

### 2. 学情分析
（分析学生基础和学习特点）

## 三、教学重点难点

### 1. 教学重点
（列出2-3个核心重点）

### 2. 教学难点
（列出1-2个主要难点及突破策略）

## 四、教学方法
（列举使用的教学方法，如：讲授法、讨论法、案例教学法等）

## 五、教学流程

### （一）导入新课（5分钟）
（设计引人入胜的导入环节）

### （二）新知讲授（20分钟）
（详细的教学过程设计，包含师生互动）

### （三）课堂练习（10分钟）
（设计针对性练习活动）

### （四）课堂小结（5分钟）
（总结归纳本节课内容）

### （五）作业布置（5分钟）
（布置分层作业）

## 六、板书设计
（清晰的板书结构）

## 七、教学反思预设
（可能遇到的问题及应对策略）

重要提示：
- 如果用户提供了教材内容，必须基于该教材内容来设计教案
- 教案中的知识点、案例、例题都要从教材中提取或延伸
- 内容详实具体，可直接用于课堂教学
- 注重学生主体性和互动性
- 体现现代教育理念
- 字数在2000-3000字之间'''
        },
        {
            'role': 'user',
            'content': user_message
        }
    ]
    
    def generate_stream():
        full_response = ""
        for chunk in stream_ai_api(messages, max_tokens=3000):
            yield chunk
            try:
                chunk_data = json.loads(chunk.replace('data: ', '').strip())
                if 'choices' in chunk_data:
                    content = chunk_data['choices'][0].get('delta', {}).get('content', '')
                    full_response += content
                elif 'error' in chunk_data:
                    full_response = chunk_data['error']
            except:
                pass
    
    return Response(
        generate_stream(),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@ai_bp.route('/extract-knowledge-points', methods=['POST'])
def extract_knowledge_points_api():
    """提取重难知识点"""
    data = request.get_json()
    lesson_plan = data.get('lesson_plan', '')
    
    if not lesson_plan:
        return jsonify({'error': '缺少教案内容'})
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的教育知识分析专家。请从教案中智能提取重难知识点，并以JSON格式返回结果。

请按以下格式返回JSON数组：
[
  {
    "title": "知识点名称",
    "content": "知识点的详细说明和解析",
    "type": "key"  // key=重点, difficult=难点, normal=一般
  },
  ...
]

要求：
1. 提取5-8个重要知识点
2. 重点（key）：必须掌握的核心概念和原理
3. 难点（difficult）：学生容易混淆或难以理解的内容
4. 一般（normal）：需要了解的基础性内容
5. 每个知识点的说明要清晰易懂
6. 确保JSON格式正确'''
        },
        {
            'role': 'user',
            'content': f'请从以下教案中提取重难知识点：\n\n{lesson_plan}'
        }
    ]
    
    response_text = call_ai_api(messages, max_tokens=1500)
    
    # 尝试解析JSON
    try:
        # 清理可能的markdown标记
        clean_response = response_text.strip()
        if clean_response.startswith('```json'):
            clean_response = clean_response[7:]
        if clean_response.startswith('```'):
            clean_response = clean_response[3:]
        if clean_response.endswith('```'):
            clean_response = clean_response[:-3]
        clean_response = clean_response.strip()
        
        knowledge_points = json.loads(clean_response)
        
        # 验证数据结构
        if isinstance(knowledge_points, list):
            for point in knowledge_points:
                if 'title' not in point:
                    point['title'] = '未命名知识点'
                if 'content' not in point:
                    point['content'] = '暂无详细说明'
                if 'type' not in point:
                    point['type'] = 'normal'
            
            return jsonify({'knowledge_points': knowledge_points})
        else:
            raise ValueError("Invalid format")
            
    except (json.JSONDecodeError, ValueError) as e:
        # 如果JSON解析失败，返回格式化的文本
        return jsonify({
            'knowledge_points': [
                {
                    'title': '知识点提取',
                    'content': response_text,
                    'type': 'key'
                }
            ],
            'raw_text': response_text
        })

@ai_bp.route('/extract-knowledge-points-stream', methods=['POST'])
def extract_knowledge_points_stream():
    """流式提取重难知识点"""
    data = request.get_json()
    lesson_plan = data.get('lesson_plan', '')
    
    if not lesson_plan:
        return jsonify({'error': '缺少教案内容'})
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的教育知识分析专家。请从教案中智能提取重难知识点。

请按以下格式输出，每个知识点用清晰的标题和分类标识：

## ★ 知识点 1：[知识点名称]
**类型：[重点/难点/一般]**

[知识点的详细说明和解析，包含具体的教学要点]

---

## ★ 知识点 2：[知识点名称]
**类型：[重点/难点/一般]**

[详细说明...]

---

要求：
1. 提取5-8个重要知识点
2. 重点：必须掌握的核心概念和原理（标记为"重点"）
3. 难点：学生容易混淆或难以理解的内容（标记为"难点"）
4. 一般：需要了解的基础性内容（标记为"一般"）
5. 每个知识点的说明要清晰易懂
6. 使用Markdown格式输出'''
        },
        {
            'role': 'user',
            'content': f'请从以下教案中提取重难知识点：\n\n{lesson_plan}'
        }
    ]
    
    def generate_stream():
        for chunk in stream_ai_api(messages, max_tokens=2000):
            yield chunk
    
    return Response(
        generate_stream(),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@ai_bp.route('/generate-quiz', methods=['POST'])
def generate_quiz_api():
    """智能出题"""
    data = request.get_json()
    knowledge_points = data.get('knowledge_points', '')
    
    if not knowledge_points:
        return jsonify({'error': '缺少知识点内容'})
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的命题专家。请根据知识点生成一套完整的练习题。

请按照以下HTML格式返回题目内容：

<div class="knowledge-point-badge">知识点编号</div>
<h2 class="quiz-title">知识点名称</h2>

<div class="question-section">
  <h3 class="question-type-title">题型一：选择题（基础题8道）</h3>
  <p class="teaching-goal"><strong>教学目标：</strong>检测学生对基础知识的记忆与理解</p>
  
  <div class="question-item">
    <p class="question-text"><strong>1.</strong> 题目内容</p>
    <div class="options">
      <p>A. 选项A &nbsp;&nbsp; B. 选项B &nbsp;&nbsp; C. 选项C &nbsp;&nbsp; D. 选项D</p>
    </div>
  </div>
  
  <!-- 更多选择题 -->
</div>

<div class="question-section">
  <h3 class="question-type-title">题型二：填空题（提高题6道）</h3>
  <p class="teaching-goal"><strong>教学目标：</strong>检测学生对关键概念的准确掌握</p>
  
  <div class="question-item">
    <p class="question-text"><strong>1.</strong> 题目内容______</p>
  </div>
  
  <!-- 更多填空题 -->
</div>

<div class="question-section">
  <h3 class="question-type-title">题型三：简答题（综合题6道）</h3>
  <p class="teaching-goal"><strong>教学目标：</strong>检测学生的综合应用和分析能力</p>
  
  <div class="question-item">
    <p class="question-text"><strong>1.</strong> 题目内容</p>
  </div>
  
  <!-- 更多简答题 -->
</div>

要求：
1. 总共20道题（选择8道+填空6道+简答6道）
2. 基础题占60%，提高题占40%
3. 题目难度梯度合理
4. 涵盖所有重要知识点
5. 题目表述清晰准确
6. 确保HTML标签正确闭合'''
        },
        {
            'role': 'user',
            'content': f'请根据以下知识点生成练习题：\n\n{knowledge_points}'
        }
    ]
    
    quiz_content = call_ai_api(messages, max_tokens=2500)
    
    # 清理响应内容
    quiz_content = quiz_content.strip()
    if quiz_content.startswith('```html'):
        quiz_content = quiz_content[7:]
    if quiz_content.startswith('```'):
        quiz_content = quiz_content[3:]
    if quiz_content.endswith('```'):
        quiz_content = quiz_content[:-3]
    quiz_content = quiz_content.strip()
    
    return jsonify({
        'quiz_content': quiz_content,
        'total_questions': 20,
        'difficulty_ratio': '60%基础 + 40%提高'
    })

@ai_bp.route('/generate-quiz-stream', methods=['POST'])
def generate_quiz_stream():
    """流式生成习题（支持自定义要求）"""
    data = request.get_json()
    knowledge_points = data.get('knowledge_points', '')
    custom_requirements = data.get('custom_requirements', '')  # 用户自定义要求
    
    if not knowledge_points:
        return jsonify({'error': '缺少知识点内容'})
    
    # 构建用户消息
    user_message = f'请根据以下知识点生成练习题：\n\n{knowledge_points}\n\n'
    
    # 如果有自定义要求，添加到prompt中
    if custom_requirements and len(custom_requirements.strip()) > 0:
        user_message += f'''【用户的出题要求】
请严格按照以下用户要求生成练习题：

{custom_requirements}

---
'''
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的命题专家。请根据知识点生成一套完整的练习题。

请使用纯文本格式输出题目，不要使用任何HTML标签。

输出格式要求：

## 一、选择题（共X道）

1. 题目内容？
   A. 选项A内容
   B. 选项B内容
   C. 选项C内容
   D. 选项D内容

2. 题目内容？
   A. ...
   ...

## 二、填空题（共X道）

1. 题目内容____________。
2. ...

## 三、简答题（共X道）

1. 题目内容
   （答题要点提示）

2. ...

重要提示：
- 使用纯文本格式，不要包含HTML标签
- 如果用户有特定的出题要求（数量、难度、题型等），必须严格按照用户的要求执行
- 题目表述清晰准确
- 每道题之间用空行分隔
- 选项使用A. B. C. D.格式'''
        },
        {
            'role': 'user',
            'content': user_message
        }
    ]
    
    def generate_stream():
        for chunk in stream_ai_api(messages, max_tokens=3000):
            yield chunk
    
    return Response(
        generate_stream(),
        mimetype='text/event-stream',
        headers={
            'Cache-Control': 'no-cache',
            'X-Accel-Buffering': 'no',
            'Connection': 'keep-alive'
        }
    )

@ai_bp.route('/generate-cards', methods=['POST'])
def generate_cards_api():
    """生成复习卡片"""
    data = request.get_json()
    knowledge_points = data.get('knowledge_points', '')
    
    if not knowledge_points:
        return jsonify({'error': '缺少知识点内容'})
    
    messages = [
        {
            'role': 'system',
            'content': '''你是一个专业的学习卡片制作专家。请将知识点转化为易于记忆的复习卡片。

请以JSON数组格式返回卡片数据：
[
  {
    "icon": "bi-book",  // Bootstrap Icons图标名
    "title": "卡片标题",
    "content": "卡片内容（简洁精炼，突出重点）",
    "tag": "分类标签"
  },
  ...
]

图标建议：
- 核心概念用: bi-lightbulb, bi-star, bi-book
- 重点公式用: bi-calculator, bi-graph-up
- 方法技巧用: bi-tools, bi-lightning
- 易错点用: bi-exclamation-triangle, bi-shield-exclamation
- 应用实例用: bi-puzzle, bi-diagram-3
- 总结归纳用: bi-check-circle, bi-journal-text

要求：
1. 生成6-10张卡片
2. 每张卡片聚焦一个核心点
3. 内容简洁明了，便于快速记忆
4. 使用项目符号或编号组织信息
5. 包含关键词高亮提示
6. 确保JSON格式正确'''
        },
        {
            'role': 'user',
            'content': f'请将以下知识点转化为复习卡片：\n\n{knowledge_points}'
        }
    ]
    
    response_text = call_ai_api(messages, max_tokens=1500)
    
    # 解析JSON
    try:
        clean_response = response_text.strip()
        if clean_response.startswith('```json'):
            clean_response = clean_response[7:]
        if clean_response.startswith('```'):
            clean_response = clean_response[3:]
        if clean_response.endswith('```'):
            clean_response = clean_response[:-3]
        clean_response = clean_response.strip()
        
        cards = json.loads(clean_response)
        
        # 验证并补充默认值
        if isinstance(cards, list):
            for card in cards:
                if 'icon' not in card:
                    card['icon'] = 'bi-book'
                if 'title' not in card:
                    card['title'] = '未命名卡片'
                if 'content' not in card:
                    card['content'] = '暂无内容'
                if 'tag' not in card:
                    card['tag'] = '知识点'
            
            return jsonify({'cards': cards})
        else:
            raise ValueError("Invalid format")
            
    except (json.JSONDecodeError, ValueError) as e:
        # JSON解析失败时返回文本格式
        return jsonify({
            'cards': [
                {
                    'icon': 'bi-journal-text',
                    'title': '复习要点',
                    'content': response_text,
                    'tag': '总结'
                }
            ],
            'raw_text': response_text
        })

# 增强版对话接口 - 支持上下文修改
@ai_bp.route('/chat/enhanced', methods=['POST'])
def chat_enhanced():
    """增强版对话接口，支持修改各类内容"""
    data = request.get_json()
    message = data.get('message', '')
    context = data.get('context', 'lesson_plan')  # lesson_plan, knowledge_points, quiz
    current_content = data.get('current_content', '')
    
    context_prompts = {
        'lesson_plan': {
            'system': '你是一个专业的教学设计助手。用户会提出对教案的修改意见，你需要根据意见修改并返回完整的修改后的教案内容。保持原有的格式结构，只修改用户要求的部分。',
            'prefix': '以下是当前教案内容，请根据用户的修改意见进行调整：\n\n'
        },
        'knowledge_points': {
            'system': '你是一个专业的教育知识管理助手。用户会对知识点提出增删改的要求，你需要按要求调整知识点列表，以JSON数组格式返回修改后的知识点。',
            'prefix': '以下是当前的知识点列表，请根据用户要求进行调整：\n\n'
        },
        'quiz': {
            'system': '你是一个专业的命题助手。用户会对练习题提出修改要求（如数量、难度、题型等），你需要按要求重新生成题目，保持原有的HTML格式。',
            'prefix': '以下是当前的练习题内容，请根据用户要求进行调整：\n\n'
        }
    }
    
    prompt_config = context_prompts.get(context, context_prompts['lesson_plan'])
    
    messages = [
        {'role': 'system', 'content': prompt_config['system']},
        {'role': 'user', 'content': f"{prompt_config['prefix']}{current_content}\n\n用户修改要求：{message}\n\n请返回修改后的完整内容："}
    ]
    
    reply = call_ai_api(messages, max_tokens=2000)
    
    result = {'reply': f'已根据您的要求完成调整：{message}'}
    
    # 根据上下文类型返回不同的更新字段
    if context == 'lesson_plan':
        result['updated_content'] = reply
    elif context == 'knowledge_points':
        try:
            clean_reply = reply.strip()
            if clean_reply.startswith('```json'):
                clean_reply = clean_reply[7:]
            if clean_reply.startswith('```'):
                clean_reply = clean_reply[3:]
            if clean_reply.endswith('```'):
                clean_reply = clean_reply[:-3]
            result['updated_knowledge_points'] = json.loads(clean_reply.strip())
        except:
            result['updated_knowledge_points'] = [{'title': '更新后的知识点', 'content': reply, 'type': 'key'}]
    elif context == 'quiz':
        clean_reply = reply.strip()
        if clean_reply.startswith('```html'):
            clean_reply = clean_reply[7:]
        if clean_reply.startswith('```'):
            clean_reply = clean_reply[3:]
        if clean_reply.endswith('```'):
            clean_reply = clean_reply[:-3]
        result['updated_quiz'] = clean_reply.strip()
    
    return jsonify(result)

# ========== 进度续接与节奏调整模块相关接口 ==========

@ai_bp.route('/analyze-comprehension', methods=['POST'])
def analyze_comprehension():
    """分析学生理解程度"""
    from models import ContentNode, LearningProgress, ComprehensionRecord, db
    from datetime import datetime
    
    data = request.get_json()
    user_id = data.get('user_id')
    content_node_id = data.get('content_node_id')
    answer = data.get('answer')
    
    if not user_id or not content_node_id or not answer:
        return jsonify({'error': '缺少必要参数'})
    
    # 获取内容节点
    content_node = ContentNode.query.get(content_node_id)
    if not content_node:
        return jsonify({'error': '内容节点不存在'})
    
    # 构建分析提示
    messages = [
        {
            'role': 'system',
            'content': '''你是一个教育评估专家，负责分析学生对知识点的理解程度。
请根据学生的回答和知识点内容，评估学生的理解程度（0-100分），并提供详细的反馈。

评估标准：
- 0-30分：基本不理解，回答与知识点无关
- 31-60分：部分理解，有一定的认识但存在错误
- 61-80分：大部分理解，基本正确但不够深入
- 81-100分：完全理解，回答全面准确

请返回以下格式：
{
  "score": 85,
  "feedback": "你的回答很好，准确理解了知识点的核心内容...",
  "missing_points": ["知识点1", "知识点2"]
}'''
        },
        {
            'role': 'user',
            'content': f'''知识点内容：{content_node.content}

学生回答：{answer}

请评估学生的理解程度并提供反馈。'''
        }
    ]
    
    # 调用AI API分析
    response = call_ai_api(messages, max_tokens=500)
    
    # 解析结果
    try:
        result = json.loads(response)
        score = result.get('score', 0)
        feedback = result.get('feedback', '')
        missing_points = result.get('missing_points', [])
        
        # 保存理解记录
        record = ComprehensionRecord(
            user_id=user_id,
            content_node_id=content_node_id,
            question=content_node.title,
            answer=answer,
            score=score,
            feedback=feedback
        )
        db.session.add(record)
        
        # 更新学习进度
        progress = LearningProgress.query.filter_by(
            user_id=user_id,
            content_node_id=content_node_id
        ).first()
        
        if progress:
            progress.comprehension_level = score
            progress.last_accessed = datetime.utcnow()
        else:
            progress = LearningProgress(
                user_id=user_id,
                content_node_id=content_node_id,
                comprehension_level=score,
                last_accessed=datetime.utcnow()
            )
            db.session.add(progress)
        
        db.session.commit()
        
        return jsonify({
            'success': True,
            'score': score,
            'feedback': feedback,
            'missing_points': missing_points,
            'progress': {
                'node_id': content_node_id,
                'comprehension_level': score,
                'last_accessed': progress.last_accessed.isoformat()
            }
        })
    except Exception as e:
        return jsonify({'error': f'分析失败：{str(e)}'})

@ai_bp.route('/get-learning-progress', methods=['POST'])
def get_learning_progress():
    """获取学习进度"""
    from models import LearningProgress
    
    data = request.get_json()
    user_id = data.get('user_id')
    
    if not user_id:
        return jsonify({'error': '缺少用户ID'})
    
    # 获取用户的学习进度
    progress_records = LearningProgress.query.filter_by(user_id=user_id).all()
    
    # 构建进度数据
    progress_data = []
    for record in progress_records:
        node = record.content_node
        progress_data.append({
            'node_id': node.id,
            'title': node.title,
            'level': node.level,
            'order_index': node.order_index,
            'completed': record.completed,
            'comprehension_level': record.comprehension_level,
            'last_accessed': record.last_accessed.isoformat(),
            'time_spent': record.time_spent
        })
    
    return jsonify({'progress': progress_data})

@ai_bp.route('/locate-difficult-nodes', methods=['POST'])
def locate_difficult_nodes():
    """定位未理解的节点"""
    from models import LearningProgress
    
    data = request.get_json()
    user_id = data.get('user_id')
    
    if not user_id:
        return jsonify({'error': '缺少用户ID'})
    
    # 获取理解程度低于60的节点
    difficult_nodes = LearningProgress.query.filter(
        LearningProgress.user_id == user_id,
        LearningProgress.comprehension_level < 60
    ).all()
    
    # 构建困难节点数据
    difficult_data = []
    for record in difficult_nodes:
        node = record.content_node
        difficult_data.append({
            'node_id': node.id,
            'title': node.title,
            'comprehension_level': record.comprehension_level,
            'last_accessed': record.last_accessed.isoformat(),
            'time_spent': record.time_spent
        })
    
    return jsonify({'difficult_nodes': difficult_data})

@ai_bp.route('/update-progress', methods=['POST'])
def update_progress():
    """更新学习进度"""
    from models import LearningProgress, db
    from datetime import datetime
    
    data = request.get_json()
    user_id = data.get('user_id')
    content_node_id = data.get('content_node_id')
    completed = data.get('completed', False)
    time_spent = data.get('time_spent', 0)
    
    if not user_id or not content_node_id:
        return jsonify({'error': '缺少必要参数'})
    
    # 查找或创建进度记录
    progress = LearningProgress.query.filter_by(
        user_id=user_id,
        content_node_id=content_node_id
    ).first()
    
    if progress:
        progress.completed = completed
        progress.time_spent += time_spent
        progress.last_accessed = datetime.utcnow()
    else:
        progress = LearningProgress(
            user_id=user_id,
            content_node_id=content_node_id,
            completed=completed,
            time_spent=time_spent,
            last_accessed=datetime.utcnow()
        )
        db.session.add(progress)
    
    db.session.commit()
    
    return jsonify({
        'success': True,
        'progress': {
            'node_id': content_node_id,
            'completed': completed,
            'time_spent': progress.time_spent,
            'last_accessed': progress.last_accessed.isoformat()
        }
    })

@ai_bp.route('/get-content-nodes', methods=['GET'])
def get_content_nodes():
    """获取内容节点列表"""
    from models import ContentNode
    
    # 获取所有内容节点
    nodes = ContentNode.query.order_by(ContentNode.level, ContentNode.order_index).all()
    
    # 构建节点数据
    node_data = []
    for node in nodes:
        node_data.append({
            'id': node.id,
            'title': node.title,
            'content': node.content,
            'parent_id': node.parent_id,
            'level': node.level,
            'order_index': node.order_index
        })
    
    return jsonify({'nodes': node_data})

@ai_bp.route('/create-content-node', methods=['POST'])
def create_content_node():
    """创建内容节点"""
    from models import ContentNode, db
    
    data = request.get_json()
    title = data.get('title')
    content = data.get('content')
    parent_id = data.get('parent_id')
    level = data.get('level', 1)
    order_index = data.get('order_index', 0)
    
    if not title:
        return jsonify({'error': '缺少标题'})
    
    # 创建新节点
    node = ContentNode(
        title=title,
        content=content,
        parent_id=parent_id,
        level=level,
        order_index=order_index
    )
    
    db.session.add(node)
    db.session.commit()
    
    return jsonify({
        'success': True,
        'node': {
            'id': node.id,
            'title': node.title,
            'content': node.content,
            'parent_id': node.parent_id,
            'level': node.level,
            'order_index': node.order_index
        }
    })

@ai_bp.route('/delete-content-node', methods=['POST'])
def delete_content_node():
    """删除内容节点"""
    from models import ContentNode, db
    
    data = request.get_json()
    node_id = data.get('node_id')
    
    if not node_id:
        return jsonify({'error': '缺少节点ID'})
    
    # 查找并删除节点
    node = ContentNode.query.get(node_id)
    if not node:
        return jsonify({'error': '节点不存在'})
    
    db.session.delete(node)
    db.session.commit()
    
    return jsonify({'success': True})

@ai_bp.route('/update-content-node', methods=['POST'])
def update_content_node():
    """更新内容节点"""
    from models import ContentNode, db
    
    data = request.get_json()
    node_id = data.get('node_id')
    title = data.get('title')
    content = data.get('content')
    parent_id = data.get('parent_id')
    level = data.get('level')
    order_index = data.get('order_index')
    
    if not node_id:
        return jsonify({'error': '缺少节点ID'})
    
    # 查找并更新节点
    node = ContentNode.query.get(node_id)
    if not node:
        return jsonify({'error': '节点不存在'})
    
    if title:
        node.title = title
    if content is not None:
        node.content = content
    if parent_id is not None:
        node.parent_id = parent_id
    if level is not None:
        node.level = level
    if order_index is not None:
        node.order_index = order_index
    
    db.session.commit()
    
    return jsonify({
        'success': True,
        'node': {
            'id': node.id,
            'title': node.title,
            'content': node.content,
            'parent_id': node.parent_id,
            'level': node.level,
            'order_index': node.order_index
        }
    })

@ai_bp.route('/ensure-user', methods=['POST'])
def ensure_user():
    """确保用户存在"""
    from models import User, db
    from werkzeug.security import generate_password_hash
    
    data = request.get_json()
    user_id = data.get('user_id')
    
    if not user_id:
        return jsonify({'error': '缺少用户ID'})
    
    # 查找用户
    user = User.query.get(user_id)
    
    # 如果用户不存在，创建一个新用户
    if not user:
        user = User(
            id=user_id,
            username=f'user{user_id}',
            email=f'user{user_id}@example.com',
            password_hash=generate_password_hash('password123', method='pbkdf2:sha256'),
            is_admin=False
        )
        db.session.add(user)
        db.session.commit()
    
    return jsonify({
        'success': True,
        'user': {
            'id': user.id,
            'username': user.username,
            'email': user.email
        }
    })
